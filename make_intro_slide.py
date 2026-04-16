#!/usr/bin/env python3
# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx"]
# ///
"""
One-page high-level architecture slide for Trina-Pi UP201/301.
5 layers, left-to-right flow:
  AI Agent  →  Skills  →  COM Port / UART  →  Trina-Pi  →  Firmware + Shell
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x11, 0x17)
TITLE_BG  = RGBColor(0x16, 0x1B, 0x22)

C_AGENT   = RGBColor(0xFF, 0xA6, 0x57)   # orange   – AI Agent
C_SKILL   = RGBColor(0x3F, 0xB9, 0x50)   # green    – Skills
C_UART    = RGBColor(0xFF, 0xC7, 0x3B)   # yellow   – COM / UART
C_HW      = RGBColor(0x58, 0xA6, 0xFF)   # blue     – Trina-Pi
C_FW      = RGBColor(0xBC, 0x8C, 0xFF)   # purple   – Firmware / Shell

F_AGENT   = RGBColor(0x22, 0x12, 0x06)
F_SKILL   = RGBColor(0x0A, 0x22, 0x16)
F_UART    = RGBColor(0x22, 0x1A, 0x04)
F_HW      = RGBColor(0x08, 0x1A, 0x32)
F_FW      = RGBColor(0x18, 0x10, 0x2C)

WHITE     = RGBColor(0xE6, 0xED, 0xF3)
LGRAY     = RGBColor(0xC9, 0xD1, 0xD9)
GRAY      = RGBColor(0x8B, 0x94, 0x9E)
DGRAY     = RGBColor(0x28, 0x32, 0x3E)

# ── Helpers ───────────────────────────────────────────────────────────────────

def box(slide, x, y, w, h, fill, brd=None, brd_pt=2.0):
    sp = slide.shapes.add_shape(
        5, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if brd:
        sp.line.color.rgb = brd
        sp.line.width = Pt(brd_pt)
    else:
        sp.line.fill.background()
    sp.text_frame.text = ""
    return sp


def txt(slide, x, y, w, h, lines):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(2)
        r = p.add_run()
        r.text    = text
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = color


def arrow_zone(slide, x, y, w, h, top_label, bot_label, color):
    """Vertical arrow zone between two blocks."""
    # central chevron
    txt(slide, x, y + h/2 - 0.22, w, 0.42, [
        ("→", 22, True, color, PP_ALIGN.CENTER),
    ])
    if top_label:
        txt(slide, x, y + h/2 - 0.55, w, 0.28, [
            (top_label, 7.5, False, LGRAY, PP_ALIGN.CENTER),
        ])
    if bot_label:
        txt(slide, x, y + h/2 + 0.22, w, 0.28, [
            (bot_label, 7.5, False, GRAY, PP_ALIGN.CENTER),
        ])


# ── Slide ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG


# ══════════════════════════════════════════════════════════════════════════════
# TITLE
# ══════════════════════════════════════════════════════════════════════════════
box(slide, 0, 0, 13.33, 0.88, TITLE_BG)
txt(slide, 0.3, 0.05, 12.7, 0.44, [
    ("Trina-Pi  UP201/301  ×  AI Agent Platform",
     23, True, WHITE, PP_ALIGN.CENTER),
])
txt(slide, 0.3, 0.50, 12.7, 0.30, [
    ("Edge Audio Intelligence  ·  Bone Microphone  ·  UART  ·  Custom Skills  ·  Claude Code  ·  Gemini CLI",
     10, False, GRAY, PP_ALIGN.CENTER),
])


# ══════════════════════════════════════════════════════════════════════════════
# LAYOUT  —  5 columns + 4 arrow gaps
#
#  col widths ("):  2.55  0.55  2.30  0.55  1.60  0.55  1.80  0.55  2.48
#  x offsets  ("):  0.20        3.30        6.15        8.30        10.65
# ══════════════════════════════════════════════════════════════════════════════
BLK_Y = 1.02
BLK_H = 6.0

# x coords
X_AGENT = 0.20;  W_AGENT = 2.55
X_ARR1  = 2.75;  W_ARR   = 0.55
X_SKILL = 3.30;  W_SKILL = 2.30
X_ARR2  = 5.60;  # same W_ARR
X_UART  = 6.15;  W_UART  = 1.90
X_ARR3  = 8.05;  # same W_ARR
X_HW    = 8.60;  W_HW    = 1.70
X_ARR4  = 10.30; # same W_ARR
X_FW    = 10.85; W_FW    = 2.28

# ── AI AGENT block ────────────────────────────────────────────────────────────
box(slide, X_AGENT, BLK_Y, W_AGENT, BLK_H, F_AGENT, C_AGENT, 2.0)

txt(slide, X_AGENT+0.10, BLK_Y+0.15, W_AGENT-0.2, 0.30, [
    ("AI AGENT", 10, True, C_AGENT, PP_ALIGN.CENTER),
])
# divider
box(slide, X_AGENT+0.15, BLK_Y+0.52, W_AGENT-0.30, 0.02, RGBColor(0x3A,0x28,0x12))

# Claude
txt(slide, X_AGENT+0.10, BLK_Y+0.70, W_AGENT-0.2, 0.40, [
    ("Claude Code", 15, True, C_AGENT, PP_ALIGN.CENTER),
])
txt(slide, X_AGENT+0.10, BLK_Y+1.10, W_AGENT-0.2, 0.26, [
    ("Anthropic", 9, False, LGRAY, PP_ALIGN.CENTER),
])

# divider
box(slide, X_AGENT+0.15, BLK_Y+2.0, W_AGENT-0.30, 0.02, RGBColor(0x3A,0x28,0x12))

# Gemini
txt(slide, X_AGENT+0.10, BLK_Y+2.18, W_AGENT-0.2, 0.40, [
    ("Gemini CLI", 15, True, C_AGENT, PP_ALIGN.CENTER),
])
txt(slide, X_AGENT+0.10, BLK_Y+2.58, W_AGENT-0.2, 0.26, [
    ("Google", 9, False, LGRAY, PP_ALIGN.CENTER),
])

# divider
box(slide, X_AGENT+0.15, BLK_Y+3.50, W_AGENT-0.30, 0.02, RGBColor(0x3A,0x28,0x12))

txt(slide, X_AGENT+0.10, BLK_Y+3.65, W_AGENT-0.2, 0.26, [
    ("reads device state", 8.5, False, GRAY, PP_ALIGN.CENTER),
])
txt(slide, X_AGENT+0.10, BLK_Y+3.95, W_AGENT-0.2, 0.26, [
    ("sends commands", 8.5, False, GRAY, PP_ALIGN.CENTER),
])
txt(slide, X_AGENT+0.10, BLK_Y+4.25, W_AGENT-0.2, 0.26, [
    ("processes audio", 8.5, False, GRAY, PP_ALIGN.CENTER),
])


# ── ARROW 1 ───────────────────────────────────────────────────────────────────
arrow_zone(slide, X_ARR1, BLK_Y, W_ARR, BLK_H, "invoke", "result", C_AGENT)


# ── SKILLS block ──────────────────────────────────────────────────────────────
box(slide, X_SKILL, BLK_Y, W_SKILL, BLK_H, F_SKILL, C_SKILL, 2.0)

txt(slide, X_SKILL+0.10, BLK_Y+0.15, W_SKILL-0.2, 0.30, [
    ("SKILLS", 10, True, C_SKILL, PP_ALIGN.CENTER),
])
box(slide, X_SKILL+0.15, BLK_Y+0.52, W_SKILL-0.30, 0.02, RGBColor(0x12,0x38,0x20))

# .md spec
txt(slide, X_SKILL+0.10, BLK_Y+0.68, W_SKILL-0.2, 0.30, [
    (".md  spec", 13, True, C_SKILL, PP_ALIGN.CENTER),
])
txt(slide, X_SKILL+0.10, BLK_Y+0.98, W_SKILL-0.2, 0.30, [
    ("Markdown description", 8.5, False, LGRAY, PP_ALIGN.CENTER),
])
txt(slide, X_SKILL+0.10, BLK_Y+1.26, W_SKILL-0.2, 0.30, [
    ("for the AI agent", 8.5, False, GRAY, PP_ALIGN.CENTER),
])

box(slide, X_SKILL+0.15, BLK_Y+1.70, W_SKILL-0.30, 0.02, RGBColor(0x12,0x38,0x20))

# .py script
txt(slide, X_SKILL+0.10, BLK_Y+1.85, W_SKILL-0.2, 0.30, [
    (".py  script", 13, True, C_SKILL, PP_ALIGN.CENTER),
])
txt(slide, X_SKILL+0.10, BLK_Y+2.15, W_SKILL-0.2, 0.30, [
    ("Python backend", 8.5, False, LGRAY, PP_ALIGN.CENTER),
])
txt(slide, X_SKILL+0.10, BLK_Y+2.43, W_SKILL-0.2, 0.30, [
    ("runs with  uv", 8.5, False, GRAY, PP_ALIGN.CENTER),
])

box(slide, X_SKILL+0.15, BLK_Y+2.88, W_SKILL-0.30, 0.02, RGBColor(0x12,0x38,0x20))

# skill names
txt(slide, X_SKILL+0.10, BLK_Y+3.05, W_SKILL-0.2, 0.30, [
    ("/mcu", 16, True, C_SKILL, PP_ALIGN.CENTER),
])
txt(slide, X_SKILL+0.10, BLK_Y+3.38, W_SKILL-0.2, 0.26, [
    ("MCU Shell Interface", 8, False, LGRAY, PP_ALIGN.CENTER),
])

box(slide, X_SKILL+0.15, BLK_Y+3.75, W_SKILL-0.30, 0.02, RGBColor(0x12,0x38,0x20))

txt(slide, X_SKILL+0.10, BLK_Y+3.92, W_SKILL-0.2, 0.30, [
    ("/pet", 16, True, C_SKILL, PP_ALIGN.CENTER),
])
txt(slide, X_SKILL+0.10, BLK_Y+4.25, W_SKILL-0.2, 0.26, [
    ("OpenClaw Pet", 8, False, LGRAY, PP_ALIGN.CENTER),
])


# ── ARROW 2 ───────────────────────────────────────────────────────────────────
arrow_zone(slide, X_ARR2, BLK_Y, W_ARR, BLK_H, "open port", "read / write", C_SKILL)


# ── COM PORT / UART block ─────────────────────────────────────────────────────
box(slide, X_UART, BLK_Y, W_UART, BLK_H, F_UART, C_UART, 2.0)

txt(slide, X_UART+0.08, BLK_Y+0.15, W_UART-0.16, 0.30, [
    ("COM PORT", 10, True, C_UART, PP_ALIGN.CENTER),
])
box(slide, X_UART+0.12, BLK_Y+0.52, W_UART-0.24, 0.02, RGBColor(0x38,0x2A,0x08))

txt(slide, X_UART+0.08, BLK_Y+0.68, W_UART-0.16, 0.36, [
    ("UART", 18, True, C_UART, PP_ALIGN.CENTER),
])

box(slide, X_UART+0.12, BLK_Y+1.18, W_UART-0.24, 0.02, RGBColor(0x38,0x2A,0x08))

txt(slide, X_UART+0.08, BLK_Y+1.32, W_UART-0.16, 0.30, [
    ("text commands", 9, True, WHITE, PP_ALIGN.CENTER),
])
txt(slide, X_UART+0.08, BLK_Y+1.60, W_UART-0.16, 0.26, [
    ("↑  ↓", 11, False, C_UART, PP_ALIGN.CENTER),
])

box(slide, X_UART+0.12, BLK_Y+2.05, W_UART-0.24, 0.02, RGBColor(0x38,0x2A,0x08))

txt(slide, X_UART+0.08, BLK_Y+2.20, W_UART-0.16, 0.30, [
    ("PCM audio", 9, True, WHITE, PP_ALIGN.CENTER),
])
txt(slide, X_UART+0.08, BLK_Y+2.48, W_UART-0.16, 0.26, [
    ("↑  ↓", 11, False, C_UART, PP_ALIGN.CENTER),
])

box(slide, X_UART+0.12, BLK_Y+2.93, W_UART-0.24, 0.02, RGBColor(0x38,0x2A,0x08))

txt(slide, X_UART+0.08, BLK_Y+3.08, W_UART-0.16, 0.30, [
    ("binary blocks", 8.5, False, LGRAY, PP_ALIGN.CENTER),
])
txt(slide, X_UART+0.08, BLK_Y+3.36, W_UART-0.16, 0.30, [
    ("10 ms per frame", 8.5, False, GRAY, PP_ALIGN.CENTER),
])


# ── ARROW 3 ───────────────────────────────────────────────────────────────────
arrow_zone(slide, X_ARR3, BLK_Y, W_ARR, BLK_H, "cmd / data", "response", C_UART)


# ── TRINA-PI block ────────────────────────────────────────────────────────────
box(slide, X_HW, BLK_Y, W_HW, BLK_H, F_HW, C_HW, 2.0)

txt(slide, X_HW+0.08, BLK_Y+0.15, W_HW-0.16, 0.30, [
    ("TRINA-PI", 10, True, C_HW, PP_ALIGN.CENTER),
])
box(slide, X_HW+0.12, BLK_Y+0.52, W_HW-0.24, 0.02, RGBColor(0x12,0x2C,0x4A))

txt(slide, X_HW+0.08, BLK_Y+0.68, W_HW-0.16, 0.38, [
    ("UP201", 17, True, C_HW, PP_ALIGN.CENTER),
])
txt(slide, X_HW+0.08, BLK_Y+1.05, W_HW-0.16, 0.28, [
    ("UP301", 11, False, LGRAY, PP_ALIGN.CENTER),
])

box(slide, X_HW+0.12, BLK_Y+1.48, W_HW-0.24, 0.02, RGBColor(0x12,0x2C,0x4A))

txt(slide, X_HW+0.08, BLK_Y+1.62, W_HW-0.16, 0.30, [
    ("RISC-V", 12, True, C_HW, PP_ALIGN.CENTER),
])
txt(slide, X_HW+0.08, BLK_Y+1.90, W_HW-0.16, 0.26, [
    ("MCU", 9, False, LGRAY, PP_ALIGN.CENTER),
])

box(slide, X_HW+0.12, BLK_Y+2.32, W_HW-0.24, 0.02, RGBColor(0x12,0x2C,0x4A))

txt(slide, X_HW+0.08, BLK_Y+2.48, W_HW-0.16, 0.30, [
    ("Bone Mic", 10, True, WHITE, PP_ALIGN.CENTER),
])
txt(slide, X_HW+0.08, BLK_Y+2.76, W_HW-0.16, 0.26, [
    ("PDM · 16 kHz", 8.5, False, LGRAY, PP_ALIGN.CENTER),
])


# ── ARROW 4 ───────────────────────────────────────────────────────────────────
arrow_zone(slide, X_ARR4, BLK_Y, W_ARR, BLK_H, "shell API", "audio out", C_HW)


# ── FIRMWARE + SHELL block ───────────────────────────────────────────────────
box(slide, X_FW, BLK_Y, W_FW, BLK_H, F_FW, C_FW, 2.0)

txt(slide, X_FW+0.10, BLK_Y+0.15, W_FW-0.2, 0.30, [
    ("FIRMWARE", 10, True, C_FW, PP_ALIGN.CENTER),
])
box(slide, X_FW+0.15, BLK_Y+0.52, W_FW-0.30, 0.02, RGBColor(0x2A,0x1A,0x44))

txt(slide, X_FW+0.10, BLK_Y+0.68, W_FW-0.2, 0.36, [
    ("on-device", 10, False, LGRAY, PP_ALIGN.CENTER),
])
txt(slide, X_FW+0.10, BLK_Y+1.02, W_FW-0.2, 0.30, [
    ("audio pipeline", 9, False, GRAY, PP_ALIGN.CENTER),
])
txt(slide, X_FW+0.10, BLK_Y+1.30, W_FW-0.2, 0.30, [
    ("DMA capture", 9, False, GRAY, PP_ALIGN.CENTER),
])
txt(slide, X_FW+0.10, BLK_Y+1.58, W_FW-0.2, 0.30, [
    ("PCM streaming", 9, False, GRAY, PP_ALIGN.CENTER),
])

box(slide, X_FW+0.15, BLK_Y+2.05, W_FW-0.30, 0.02, RGBColor(0x2A,0x1A,0x44))

txt(slide, X_FW+0.10, BLK_Y+2.20, W_FW-0.2, 0.36, [
    ("Shell", 16, True, C_FW, PP_ALIGN.CENTER),
])
txt(slide, X_FW+0.10, BLK_Y+2.60, W_FW-0.2, 0.30, [
    ("simple command REPL", 8.5, False, LGRAY, PP_ALIGN.CENTER),
])

box(slide, X_FW+0.15, BLK_Y+3.00, W_FW-0.30, 0.02, RGBColor(0x2A,0x1A,0x44))

for i, (cmd, c) in enumerate([
    ("pet",    C_FW),
    ("stream", RGBColor(0x3F,0xB9,0x50)),
    ("record", RGBColor(0x3F,0xB9,0x50)),
    ("pdm",    C_HW),
    ("info",   C_HW),
    ("reset",  C_UART),
]):
    row, col = divmod(i, 2)
    cx = X_FW + 0.18 + col * 1.08
    cy = BLK_Y + 3.18 + row * 0.46
    sp = slide.shapes.add_shape(5, Inches(cx), Inches(cy), Inches(0.88), Inches(0.32))
    sp.fill.solid(); sp.fill.fore_color.rgb = F_FW
    sp.line.color.rgb = c; sp.line.width = Pt(1.0)
    sp.text_frame.text = ""
    p = sp.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = cmd
    r.font.size = Pt(8.5); r.font.bold = True; r.font.color.rgb = c


# ══════════════════════════════════════════════════════════════════════════════
# FOOTER
# ══════════════════════════════════════════════════════════════════════════════
box(slide, 0, 7.18, 13.33, 0.32, TITLE_BG)
txt(slide, 0.3, 7.20, 12.7, 0.28, [
    ("github.com/KaiYin77/upbeat-skills  ·  curl -fsSL …/install.sh | bash",
     8, False, GRAY, PP_ALIGN.CENTER),
])

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save("trina_pi_intro.pptx")
print("Saved: trina_pi_intro.pptx")
